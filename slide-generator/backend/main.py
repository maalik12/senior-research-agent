import os
import io
import base64
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import fitz  # PyMuPDF

app = Flask(__name__)
CORS(app)

# Simple in-memory storage for demo (use database in production)
presentations = {}

def extract_text_from_pdf(pdf_bytes):
    """Extract text content from PDF document"""
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text_content = []
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        text_content.append({
            'page': page_num + 1,
            'content': text.strip()
        })
    
    doc.close()
    return text_content

def create_presentation_slide(title, content, background_color='FFFFFF', font_size=18):
    """Create a single slide with given content"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(background_color)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame
    title_tf.paragraphs[0].font.size = Pt(32)
    title_tf.paragraphs[0].font.bold = True
    
    # Set content
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        paragraphs = content.split('\n')
        for i, paragraph in enumerate(paragraphs[:5]):  # Limit to 5 paragraphs per slide
            if paragraph.strip():
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = paragraph.strip()
                p.font.size = Pt(font_size)
                p.alignment = PP_ALIGN.LEFT
    
    return prs

@app.route('/api/analyze', methods=['POST'])
def analyze_document():
    """Analyze uploaded document and extract key points"""
    try:
        if 'document' not in request.files:
            return jsonify({'error': 'No document uploaded'}), 400
        
        file = request.files['document']
        if file.filename.endswith('.pdf'):
            pdf_bytes = file.read()
            extracted_text = extract_text_from_pdf(pdf_bytes)
            
            # Simple extraction of key points (in production, use AI)
            key_points = []
            for page in extracted_text:
                lines = page['content'].split('\n')
                for line in lines:
                    if len(line.strip()) > 20:  # Filter meaningful lines
                        key_points.append({
                            'slide_title': f"Page {page['page']} - Key Point",
                            'content': line.strip()[:200]  # Limit length
                        })
            
            return jsonify({
                'success': True,
                'extracted_points': key_points[:10],  # Limit to 10 slides
                'total_pages': len(extracted_text)
            })
        else:
            return jsonify({'error': 'Only PDF files are supported'}), 400
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/generate', methods=['POST'])
def generate_presentation():
    """Generate PowerPoint presentation based on requirements"""
    try:
        data = request.json
        slides_data = data.get('slides', [])
        background_color = data.get('background_color', 'FFFFFF')
        font_size = data.get('font_size', 18)
        
        if not slides_data:
            return jsonify({'error': 'No slide data provided'}), 400
        
        prs = Presentation()
        
        for slide_info in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(background_color)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_info.get('title', 'Untitled Slide')
            title_tf = title_shape.text_frame
            title_tf.paragraphs[0].font.size = Pt(32)
            title_tf.paragraphs[0].font.bold = True
            
            # Set content
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                
                content = slide_info.get('content', '')
                paragraphs = content.split('\n')
                for i, paragraph in enumerate(paragraphs[:6]):
                    if paragraph.strip():
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = paragraph.strip()
                        p.font.size = Pt(font_size)
                        p.alignment = PP_ALIGN.LEFT
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/templates', methods=['GET'])
def get_templates():
    """Get available presentation templates"""
    templates = [
        {
            'id': 'academic_blue',
            'name': 'Academic Blue',
            'background_color': 'E6F3FF',
            'description': 'Clean blue background suitable for academic presentations'
        },
        {
            'id': 'minimal_white',
            'name': 'Minimal White',
            'background_color': 'FFFFFF',
            'description': 'Clean white background with high contrast'
        },
        {
            'id': 'dark_professional',
            'name': 'Dark Professional',
            'background_color': '2C3E50',
            'description': 'Dark background for professional settings'
        },
        {
            'id': 'green_academic',
            'name': 'Academic Green',
            'background_color': 'E8F5E9',
            'description': 'Soft green background for educational content'
        }
    ]
    return jsonify({'templates': templates})

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)